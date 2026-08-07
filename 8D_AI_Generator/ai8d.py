"""Generate an editable 8D report from D0-D8 folders and a PowerPoint template."""

from __future__ import annotations

import argparse
import base64
import copy
import io
import json
import os
import re
import sys
import time
from pathlib import Path
from typing import Any

from json_repair import repair_json
from openai import OpenAI
import httpx
from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


API_URL_ENV = "ADAYO8D_API_URL"
API_KEY_ENV = "ADAYO8D_API_KEY"
DEFAULT_API_URL = "http://10.2.9.178:4000/v1"
DEFAULT_API_KEY = "sk-nQHNlCWBO73aAGqVwbImQfzd5NsGv4dyk4fPIAlYu1OHd79J"
MODEL = "MiniMax-M3"
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".webp"}
MAX_IMAGE_EDGE = 1600
JPEG_QUALITY = 85
MAX_IMAGES_PER_DIRECT_REQUEST = 12
IMAGES_PER_BATCH = 10
STAGE_CODES = ("d0", "d1", "d2", "d3", "d4", "d5", "d6", "d7", "d8")
STAGE_NAMES = [
    "D0 相关信息", "D1 团队成立", "D2 问题描述", "D3 临时措施", "D4 原因分析",
    "D5 长期对策", "D6 效果验证", "D7 预防措施", "D8 结案总结",
]


def get_api_url() -> str:
    """API 地址：优先读环境变量 ADAYO8D_API_URL，否则用固定默认值。"""
    return os.environ.get(API_URL_ENV, DEFAULT_API_URL)


def get_api_key() -> str:
    """API Key：优先读环境变量 ADAYO8D_API_KEY，否则用固定默认值。"""
    env_key = os.environ.get(API_KEY_ENV)
    if env_key:
        return env_key.strip()
    return DEFAULT_API_KEY


def make_openai_client(api_key: str, timeout: float) -> OpenAI:
    """创建 OpenAI 客户端，始终直连内网服务，不走系统代理。

    TUN/系统代理开启时，httpx 默认会读取系统代理并把内网请求转发到
    代理端口（导致 502），这里显式禁用代理并固定走物理网卡直连。
    """
    transport = httpx.HTTPTransport(proxy=None)
    http_client = httpx.Client(transport=transport, trust_env=False, timeout=timeout)
    return OpenAI(base_url=get_api_url(), api_key=api_key, http_client=http_client, max_retries=0)


def read_text(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "gb18030", "utf-16"):
        try:
            return raw.decode(encoding).strip()
        except UnicodeDecodeError:
            pass
    return raw.decode("utf-8", errors="replace").strip()


def collect_materials(materials_root: Path) -> dict[str, dict[str, Any]]:
    result: dict[str, dict[str, Any]] = {}
    folders = [item for item in materials_root.iterdir() if item.is_dir()]
    for stage_index, stage_name in enumerate(STAGE_NAMES):
        folder = next(
            (item for item in folders if item.name.lower().startswith(f"d{stage_index}") or item.name.startswith(f"{stage_index:02}")),
            None,
        )
        if folder is None:
            result[stage_name] = {"notes": "", "images": []}
            continue
        notes = "\n\n".join(read_text(item) for item in folder.rglob("*.txt") if item.stat().st_size)
        images = sorted(item for item in folder.rglob("*") if item.suffix.lower() in IMAGE_EXTENSIONS)
        result[stage_name] = {"notes": notes, "images": images}
    return result


def executable_folder() -> Path:
    return Path(sys.executable).resolve().parent if getattr(sys, "frozen", False) else Path(__file__).resolve().parent


def bundled_template() -> Path:
    base = Path(getattr(sys, "_MEIPASS", Path(__file__).resolve().parent))
    return base / "8D模板.pptx"


def clear_targets(path: Path) -> list[Path]:
    """List files in explicitly named D0-D8 / 00-08 material folders."""
    stage_pattern = re.compile(r"^(?:d[0-8]|0[0-8])(?:\s|$)", re.I)
    return [
        item
        for entry in path.iterdir()
        if entry.is_dir() and stage_pattern.match(entry.name)
        for item in entry.rglob("*")
        if item.is_file()
    ]


def clear_folders(path: Path) -> int:
    """Recursively remove material files while retaining the D0-D8 folders."""
    targets = clear_targets(path)
    for item in targets:
        item.unlink()
    for entry in path.iterdir():
        if entry.is_dir() and re.match(r"^(?:d[0-8]|0[0-8])(?:\s|$)", entry.name, re.I):
            for child in sorted(entry.rglob("*"), reverse=True):
                if child.is_dir() and not any(child.iterdir()):
                    child.rmdir()
    return len(targets)


def confirm_clear(path: Path) -> bool:
    count = len(clear_targets(path))
    if not count:
        print("D0-D8 文件夹中没有可清空的文件。")
        return False
    answer = input(f"将递归删除 D0-D8 文件夹中的 {count} 个文件，确认继续？(y/N): ").strip().lower()
    return answer in {"y", "yes"}


def pause(message: str = "\nPress Enter to close...") -> None:
    if getattr(sys, "frozen", False):
        try:
            input(message)
        except EOFError:
            pass


def interactive_menu() -> int:
    print("=" * 40)
    print("  Adayo 8D AI 报告生成器")
    print("  1. 生成 8D PPT 报告")
    print("  2. 清空 D0-D8 材料文件夹")
    print("  q. 退出")
    print("=" * 40)
    choice = input("请选择 (1/2/q): ").strip().lower()
    if choice == "1":
        return 1
    if choice == "2":
        return 2
    if choice == "q":
        return 0
    print("无效选项，请输入 1、2 或 q。")
    return interactive_menu()


def image_part(path: Path) -> dict[str, Any]:
    """Encode a bounded JPEG preview instead of uploading original camera files."""
    try:
        with Image.open(path) as source:
            source.thumbnail((MAX_IMAGE_EDGE, MAX_IMAGE_EDGE))
            if source.mode not in ("RGB", "L"):
                source = source.convert("RGB")
            buffer = io.BytesIO()
            source.save(buffer, format="JPEG", quality=JPEG_QUALITY, optimize=True)
    except Exception as exc:
        raise ValueError(f"无法读取图片 {path.name}: {exc}") from exc
    encoded = base64.b64encode(buffer.getvalue()).decode("ascii")
    return {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{encoded}"}}


def build_prompt(
    materials: dict[str, dict[str, Any]],
    include_images: bool = True,
    image_evidence: dict[str, list[dict[str, str]]] | None = None,
) -> list[dict[str, Any]]:
    """Build the final report request, with optional evidence text from image batches."""
    content: list[dict[str, Any]] = [{"type": "text", "text": "以下为 D0-D8 分组材料。TXT 是人员记录，图片是证据。请严格只依据材料输出。"}]
    for stage, payload in materials.items():
        content.append({"type": "text", "text": f"\n【{stage}】\n文字记录：\n{payload['notes'] or '（无）'}\n图片文件：{', '.join(image.name for image in payload['images']) or '（无）'}"})
        summaries = (image_evidence or {}).get(stage, [])
        if summaries:
            evidence_text = "\n".join(
                f"- {item['image']}：{item['summary']}" for item in summaries if item.get("summary")
            )
            if evidence_text:
                content.append({"type": "text", "text": f"图片分批排查结论（可用于报告，不要在正文复述文件名）：\n{evidence_text}"})
        if include_images:
            for image in payload["images"]:
                content.append(image_part(image))
    return content


SYSTEM_PROMPT = r"""
你是汽车电子质量 8D 报告助手。基于用户给出的分组文字和现场证据图片，输出单个 JSON 对象，不要 Markdown，不要解释。

原则：
1. 不得编造日期、数据或责任人；但应基于文字记录和图片证据完成合理的原因链、长期对策、验证方式和结案总结。没有明确日期或责任人时留空，不能用"待确认"代替内容。
2. 图片只能支持可直接观察到的现象；不能凭图片断言真实根因。
3. 使用正式、简洁的中文。每项最多两句。team 最多 6 人；countermeasures、validations 各最多 5 项。
4. evidence_summary 是每个阶段图片想证明的简短说明（不含图片文件名），没有图片则为空字符串。
5. D2/D4/D5/D6/D7/D8 的 summary 字段控制在 40 字以内（1 句），不要长篇叙述。

填写要求（关键）：
- D0 问题了解：根据材料填写 fault_date / model / trace_code / station / description；材料没有明确说明的字段留空，不要写"待确认"。
- D1 团队成立：team 列表里每人按 name / role / duty / contact / module 输出，但 role 字段保持模板给定角色（组长 / 组员），不要改写。
- D2 故障描述：customer / date / supplier / model / vehicle_model / material / quantity / failure_type / lot / location / production_date / vin / complaint_source / symptom / confirmed_symptom 全部留空字符串。
- D3 临时措施：actions 列表留空（不要生成任何条目），summary 写"待确认"。
- D4 原因分析：必须输出 occurrence_why_rows 的 Why1 到 Why5 共五行，escape_why_rows 的 Why1 到 Why3 共三行，每行含 level / problem / cause，缺一不可。即使证据不充分也必须基于材料合理推断原因链并填满全部行，不能只填一行就结束。cause 是推理出的原因答案，不能留空或写"无"。problem 和 cause 各写一到两句具体内容，不使用占位文本。
- D5 长期对策：按 action / owner / date / status 输出至少 4 条（最多 5 条）；status 默认"进行中"。action 必须覆盖具体整改动作，并自然包含与程序文件、FMEA、控制计划、SOP、经验教训库相关的关键词，以确保与 D7 预防措施表格联动。
- D6 效果验证：按 method / result / owner / date 输出。
- D7 预防措施：按 preventions 数组输出，每项含 item（程序Procedure/工作指示Work instruction/操作指示SOP/流程图Flow chart/失效模式分析D/P-FMEA/控制计划Control Plan/设计规范Design disciplines/经验教训Lessons Learned）、verdict（"是"或"否"）、comment（如为"是"则填写关联的具体对策描述，如为"否"则写"NA"）。summary 写具体一句话。
- D8 结案总结：必须分别填写 background（问题背景概述，2-3 句）、current_status（当前状态和处理结果，2-3 句）、system_strategy（系统性整改策略和措施，2-3 句）、strategy_logic（策略逻辑、横向推广和持续改进，2-3 句）。四个字段缺一不可，不能将所有内容堆在 background 中，其它字段也必须独立成段。conclusion 全文不少于 500 字。

输出键必须齐全，按以下 JSON 结构：
{
 "cover":{"title":"","prepared_by":"","reviewed_by":"","approved_by":"","date":""},
 "d0":{"fault_date":"","model":"","trace_code":"","station":"","description":"","summary":"","evidence_summary":""},
 "d1":{"team":[{"name":"","role":"","duty":"","contact":"","module":""}],"summary":"","evidence_summary":""},
 "d2":{"customer":"","date":"","supplier":"","model":"","vehicle_model":"","material":"","quantity":"","failure_type":"","lot":"","location":"","production_date":"","vin":"","complaint_source":"","symptom":"","confirmed_symptom":"","summary":"","evidence_summary":""},
 "d3":{"actions":[],"summary":"待确认","evidence_summary":""},
 "d4":{"occurrence_why_rows":[{"level":"Why1","problem":"","cause":""}],"escape_why_rows":[{"level":"Why1","problem":"","cause":""}],"root_cause":"","escape_cause":"","summary":"","evidence_summary":""},
 "d5":{"countermeasures":[{"action":"","owner":"","date":"","status":""}],"summary":"","evidence_summary":""},
 "d6":{"validations":[{"method":"","result":"","owner":"","date":""}],"summary":"","evidence_summary":""},
  "d7":{"preventions":[{"item":"","verdict":"","comment":""}],"summary":"","evidence_summary":""},
  "d8":{"background":"问题背景概述（2-3 句，必须独立填写）","current_status":"当前状态和处理结果（2-3 句，必须独立填写）","system_strategy":"系统性整改策略和措施（2-3 句，必须独立填写）","strategy_logic":"策略逻辑、横向推广和持续改进（2-3 句，必须独立填写）","conclusion":"","summary":"","evidence_summary":""}
}
""".strip()
SYSTEM_PROMPT += """

Additionally include a root-level image_pages object with d2, d4, d5, d6, d7 and d8 arrays.
Each page must contain title, summary, layout (single, two, three or four), and images (exact image file names).
Use single for screenshots, comparison images, or images that require inspection individually. Only group images that show one same activity.
Every supplied image name must appear exactly once in its own D-stage list. Do not mention any image file name in summary text.
The image-page summary is an 8D investigation report caption, not a visual description. State the investigation object, what was checked or verified, the finding or conclusion supported by the evidence, and its relevance to containment, root-cause analysis, corrective action, validation, or prevention. Do not write generic captions such as "shown in the image" or merely list visible UI elements.
当同一 D 阶段有多页时，每页的 summary 必须根据各自图片内容写出不同的小结，禁止多页使用相同文字。"""

SYSTEM_PROMPT += """

For d4, return occurrence_why_rows with exactly five ordered rows (Why1 to Why5), and escape_why_rows with exactly three ordered rows (Why1 to Why3).
Every row must contain level, problem and cause. `cause` is the answer/reason for that Why level, never a rewording of the question. The first `problem` is the original fault; later rows' `problem` should be the previous row's cause, so the chain is continuous.
The PPT has three columns: level, problem and cause. Put the answer in `cause`, not only in the `whys` list. Do not make up unsupported facts: if a level genuinely has no evidence, leave that row's cause empty.
"""


IMAGE_EVIDENCE_SYSTEM_PROMPT = """
你是汽车电子质量8D报告的现场证据分析助手。只分析本次提供的一个阶段、一批图片，输出单个 JSON 对象，不要 Markdown 或解释。
返回格式必须为：
{"observations":[{"image":"原始图片文件名","summary":"排查结论"}]}
每一张输入图片都必须且只能有一条 observation，image 必须与输入文件名完全一致。summary 使用8D排查口径，简要说明：排查/核对的对象、图片中可直接确认的现象或验证结果、它对问题界定、原因分析、措施或验证的支持。没有可确认信息时写“该图片已纳入现场证据核对，未见可独立确认的结论。”不得虚构图片中不可见的事实。summary 中不要写文件名。
""".strip()


def parse_json(text: str) -> dict[str, Any]:
    cleaned = text.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?\s*|\s*```$", "", cleaned, flags=re.S)
    match = re.search(r"\{.*\}", cleaned, flags=re.S)
    if not match:
        raise ValueError("模型没有返回 JSON 对象。")
    candidate = match.group(0)
    try:
        return json.loads(candidate)
    except json.JSONDecodeError:
        repaired = repair_json(candidate, return_objects=True)
        if not isinstance(repaired, dict):
            raise ValueError("模型返回内容无法修复为 JSON 对象。")
        return repaired


def content_warnings(data: dict[str, Any]) -> list[str]:
    """Report evidence gaps without rejecting intentionally blank factual fields."""
    d4 = data.get("d4", {})
    occurrence = d4.get("occurrence_why_rows", [])
    escape = d4.get("escape_why_rows", [])
    missing: list[str] = []
    if len(occurrence) < 5 or any(not str(row.get("cause") or "").strip() for row in occurrence[:5]):
        missing.append("D4发生原因5WHY（5行原因）")
    if len(escape) < 3 or any(not str(row.get("cause") or "").strip() for row in escape[:3]):
        missing.append("D4流出原因5WHY（3行原因）")
    if not all(str(row.get("problem") or "").strip() for row in occurrence[:5]):
        missing.append("D4发生原因problem列有空行")
    if not all(str(row.get("problem") or "").strip() for row in escape[:3]):
        missing.append("D4流出原因problem列有空行")
    if len(usable_records(data.get("d5", {}).get("countermeasures", []), "action")) < 4:
        missing.append("D5长期对策（至少4条）")
    if not usable_records(data.get("d6", {}).get("validations", []), "method"):
        missing.append("D6效果验证")
    conclusion = str(data.get("d8", {}).get("conclusion") or "").strip()
    d8 = data.get("d8", {})
    d8_fields = ("background", "current_status", "system_strategy", "strategy_logic")
    if len(re.sub(r"\s+", "", conclusion)) < 250 or any(not str(d8.get(key) or "").strip() for key in d8_fields):
        missing.append("D8结案总结（不少于250字）")
    return missing


def describe_ai_error(exc: Exception) -> str:
    """Turn common provider failures into user-actionable console messages."""
    message = str(exc)
    lowered = message.lower()
    if "401" in message or "403" in message or "api key" in lowered or "authentication" in lowered:
        return "API Key 无效或没有访问权限。请检查 ADAYO8D_API_KEY。"
    if "timeout" in lowered or "timed out" in lowered:
        return "模型响应超时。请检查网络或稍后重试。"
    if "connection" in lowered or "connect" in lowered or "dns" in lowered:
        return "无法连接模型服务。请检查 ADAYO8D_API_URL、网络和服务状态。"
    return f"模型服务返回错误：{message}"


def split_batches(items: list[Any], batch_size: int) -> list[list[Any]]:
    """Keep image requests bounded so a large report cannot exceed gateway limits."""
    if batch_size < 1:
        raise ValueError("batch_size 必须大于 0。")
    return [items[index:index + batch_size] for index in range(0, len(items), batch_size)]


def generate_image_evidence(materials: dict[str, dict[str, Any]], api_key: str) -> dict[str, list[dict[str, str]]]:
    """Analyze large image sets in bounded requests, retrying only the failed batch."""
    client = make_openai_client(api_key, 120.0)
    evidence: dict[str, list[dict[str, str]]] = {}
    for stage, payload in materials.items():
        batches = split_batches(list(payload["images"]), IMAGES_PER_BATCH)
        if not batches:
            continue
        stage_evidence: list[dict[str, str]] = []
        for batch_number, batch in enumerate(batches, start=1):
            filename_text = ", ".join(image.name for image in batch)
            batch_text = (
                f"【{stage}】第 {batch_number}/{len(batches)} 批现场证据。\n"
                f"文字记录：{payload['notes'] or '（无）'}\n"
                f"本批图片文件：{filename_text}"
            )
            messages = [
                {"role": "system", "content": IMAGE_EVIDENCE_SYSTEM_PROMPT},
                {"role": "user", "content": [{"type": "text", "text": batch_text}, *(image_part(image) for image in batch)]},
            ]
            last_error: Exception | None = None
            data: dict[str, Any] | None = None
            for attempt in range(1, 4):
                try:
                    print(
                        f"      [图像取证] {stage} 第 {batch_number}/{len(batches)} 批（{len(batch)} 张），"
                        f"第 {attempt}/3 次请求…",
                        flush=True,
                    )
                    response = client.chat.completions.create(model=MODEL, temperature=0.1, messages=messages)
                    data = parse_json(response.choices[0].message.content or "")
                    break
                except Exception as exc:
                    last_error = exc
                    print(f"      [图像取证] 本批第 {attempt}/3 次失败：{describe_ai_error(exc)}", flush=True)
                    if attempt < 3:
                        wait = 3 * attempt
                        print(f"      [图像取证] 等待 {wait} 秒后重试…", flush=True)
                        time.sleep(wait)
            if data is None:
                raise RuntimeError(
                    f"{stage} 第 {batch_number}/{len(batches)} 批图片取证失败："
                    f"{describe_ai_error(last_error or RuntimeError())}"
                )
            observations = {
                str(item.get("image") or "").strip(): str(item.get("summary") or "").strip()
                for item in data.get("observations", []) if isinstance(item, dict)
            }
            for image in batch:
                stage_evidence.append({
                    "image": image.name,
                    "summary": observations.get(image.name) or "该图片已纳入现场证据核对，未见可独立确认的结论。",
                })
        evidence[stage] = stage_evidence
    return evidence


def generate_content(materials: dict[str, dict[str, Any]], api_key: str) -> dict[str, Any]:
    client = make_openai_client(api_key, 300.0)
    started = time.monotonic()
    image_count = sum(len(payload["images"]) for payload in materials.values())
    image_evidence: dict[str, list[dict[str, str]]] = {}
    if image_count <= MAX_IMAGES_PER_DIRECT_REQUEST:
        print(f"      [模型] 共 {image_count} 张图片，使用单次图文生成模式。", flush=True)
        prompt = build_prompt(materials)
    else:
        print(
            f"      [模型] 共 {image_count} 张图片，切换为分批取证模式（每批最多 {IMAGES_PER_BATCH} 张）…",
            flush=True,
        )
        image_evidence = generate_image_evidence(materials, api_key)
        print("      [模型] 图片分批取证完成，正在根据文字和取证结论汇总生成报告…", flush=True)
        prompt = build_prompt(materials, include_images=False, image_evidence=image_evidence)
    messages = [
        {"role": "system", "content": SYSTEM_PROMPT},
        {"role": "user", "content": prompt},
    ]
    try:
        print("      [模型] 正在连接并提交材料（单次最多等待 300 秒，超时自动重试 2 次）…", flush=True)
        response = client.chat.completions.create(model=MODEL, temperature=0.1, messages=messages)
        raw = response.choices[0].message.content or ""
        print(f"      [模型] 已收到回复（{time.monotonic() - started:.1f} 秒），正在校验内容…", flush=True)
        data = parse_json(raw)
        if image_evidence:
            data["_image_evidence"] = image_evidence
        print("      [模型] 内容校验成功。", flush=True)
        return data
    except Exception as exc:
        raise RuntimeError(f"AI 调用失败（已自动重试，超时 300 秒）：{describe_ai_error(exc)}") from exc


def generate_title_from_ppt(report: Path, api_key: str) -> str:
    """Generate a concise report title from the completed PPT content."""
    presentation = Presentation(str(report))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                parts.append(shape.text.strip())
            if shape.has_table:
                parts.extend(cell.text.strip() for row in shape.table.rows for cell in row.cells if cell.text.strip())
    client = make_openai_client(api_key, 120.0)
    started = time.monotonic()
    print("      [标题] 正在发送完整报告摘要给模型…", flush=True)
    try:
        response = client.chat.completions.create(
            model=MODEL,
            temperature=0.1,
            messages=[
                {"role": "system", "content": "你是汽车质量8D报告编辑。根据完整报告内容生成正式中文标题。只输出标题，不加引号、序号或解释；不超过28个汉字。"},
                {"role": "user", "content": "\n".join(parts)[:12000]},
            ],
        )
        title = (response.choices[0].message.content or "").strip().strip('“”"')[:40]
        print(f"      [标题] 模型返回成功（{time.monotonic() - started:.1f} 秒）。", flush=True)
        return title
    except Exception as exc:
        raise RuntimeError(describe_ai_error(exc)) from exc


def locate_cover_slide(presentation: Presentation) -> int:
    """Locate the cover slide after tokens have been replaced ({{Content}} no longer exists)."""
    labels = ("编制", "审核", "批准")
    for index, slide in enumerate(presentation.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text)
            if shape.has_table:
                texts.extend(cell.text for row in shape.table.rows for cell in row.cells)
        source = re.sub(r"\s+", "", "".join(texts))
        if all(label in source for label in labels):
            return index
    raise ValueError("模板无法唯一定位 cover 页面（找不到含 编制/审核/批准 的页面）。")


def write_cover_title(report: Path, title: str) -> None:
    presentation = Presentation(str(report))
    try:
        cover_index = locate_cover_slide(presentation)
    except ValueError:
        return
    cover = presentation.slides[cover_index]
    for shape in cover.shapes:
        if shape.name == "AI_COVER_TITLE":
            shape.text = title
            break
    else:
        replace_tokens(cover.shapes[0], {"{{Content}}": title})
    presentation.save(str(report))


def text(value: Any) -> str:
    return str(value or "").strip()


def set_run_font(run: Any, size: int = 10) -> None:
    """Set Latin and East-Asian font metadata; PowerPoint otherwise falls back to SimSun."""
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    rpr = run._r.get_or_add_rPr()
    rpr.set("latin", "Microsoft YaHei")
    rpr.set("ea", "Microsoft YaHei")
    rpr.set("cs", "Microsoft YaHei")


def put_cell(table: Any, row: int, col: int, value: Any) -> None:
    cell = table.cell(row, col)
    cell.text = "" if value in (None, "", "待确认") else str(value).strip()
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            set_run_font(run)


def first_table(slide: Any) -> Any:
    table_shape = next((shape for shape in slide.shapes if shape.has_table), None)
    if table_shape is None:
        raise ValueError("模板页面未找到可填写的原生表格。")
    return table_shape


def locate_template_slides(presentation: Presentation) -> dict[str, int]:
    """Locate report pages by their visible headings instead of fragile indexes."""
    def shape_texts(shapes: Any) -> list[str]:
        result: list[str] = []
        for shape in shapes:
            if shape.has_text_frame:
                result.append(shape.text)
            if shape.has_table:
                result.extend(cell.text for row in shape.table.rows for cell in row.cells)
            if hasattr(shape, "shapes"):
                result.extend(shape_texts(shape.shapes))
        return result

    texts = [
        re.sub(r"\s+", "", "\n".join(shape_texts(slide.shapes)))
        for slide in presentation.slides
    ]
    rules = {
        "cover": (("{{Content}}",), ()),
        "d0": (("故障日期", "机型", "总成追溯条码"), ()),
        "d1": (("团队成员", "主要负责模块"), ()),
        "d2": (("客户", "供应商", "投诉来源"), ()),
        "d3": (("NO", "客户方", "在途"), ()),
        "d4_evidence": (("D4", "原因", "分析"), ("5WHY", "目录")),
        "d4_occurrence": (("D4", "5WHY", "发生"), ()),
        "d4_escape": (("D4", "5WHY", "流出"), ()),
        "d5": (("编号", "长期对策", "负责人"), ()),
        "d6": (("编号", "验证方式", "效果验证"), ()),
        "d7": (("程序", "工作指示", "控制计划"), ()),
        "d8": (("背景", "系统策略", "策略逻辑"), ()),
    }
    locations: dict[str, int] = {}
    for key, (required, forbidden) in rules.items():
        matches = [
            index for index, value in enumerate(texts)
            if all(token in value for token in required) and not any(token in value for token in forbidden)
        ]
        if len(matches) != 1:
            raise ValueError(f"模板无法唯一定位 {key} 页面（找到 {len(matches)} 页）。请检查标题文本。")
        locations[key] = matches[0]
    return locations


def replace_tokens(shape: Any, replacements: dict[str, Any]) -> None:
    if not shape.has_text_frame:
        return
    original = shape.text
    updated = original
    for token, value in replacements.items():
        updated = updated.replace(token, text(value))
    if updated != original:
        shape.text = updated
        shape.name = "AI_DYNAMIC"


def replace_content_placeholders(slide: Any, values: list[Any]) -> None:
    """Fill {{content}} paragraphs in their template order.

    A designer may put several placeholders in one text box, so a shape is not
    necessarily a single value slot.
    """
    placeholders = [
        (shape, paragraph)
        for shape in slide.shapes if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        if "{{content}}" in paragraph.text
    ]
    if len(placeholders) != len(values):
        raise ValueError(f"模板页 {{content}} 数量为 {len(placeholders)}，但程序需要填写 {len(values)} 项。")
    for (shape, paragraph), value in zip(placeholders, values):
        # Keep labels and signatures that share the paragraph with the marker
        # (for example, "编制：{{content}}").  An empty model value removes
        # only the marker and intentionally leaves the surrounding template
        # text intact.
        replacement = "" if value is None else str(value).strip()
        marker_runs = [run for run in paragraph.runs if "{{content}}" in run.text]
        if marker_runs:
            for run in marker_runs:
                run.text = run.text.replace("{{content}}", replacement)
        else:
            # Defensive fallback for a marker split across runs by a future template.
            paragraph.text = paragraph.text.replace("{{content}}", replacement)
        shape.name = "AI_DYNAMIC"


def set_paragraph_text(paragraph: Any, value: Any) -> None:
    """Replace a paragraph while retaining a deterministic Chinese font setting."""
    paragraph.text = "" if value in (None, "", "待确认") else str(value).strip()
    for run in paragraph.runs:
        set_run_font(run)


def add_text(slide: Any, value: str, left: float, top: float, width: float, height: float, size: int = 10) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.name = "AI_DYNAMIC"
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = value
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            set_run_font(run, size)


def crop_picture(slide: Any, image: Path, left: float, top: float, width: float, height: float) -> None:
    if left < 0 or top < 0 or left + width > 10.0 or top + height > 5.625:
        raise ValueError(f"image layout exceeds slide bounds: {image.name}")
    with Image.open(image) as source:
        image_ratio = source.width / source.height
    frame_ratio = width / height
    if image_ratio > frame_ratio:
        render_width = width
        render_height = width / image_ratio
        render_left = left
        render_top = top + (height - render_height) / 2
    else:
        render_height = height
        render_width = height * image_ratio
        render_left = left + (width - render_width) / 2
        render_top = top
    slide.shapes.add_picture(str(image), Inches(render_left), Inches(render_top), width=Inches(render_width), height=Inches(render_height))


def clone_slide(presentation: Presentation, source_index: int) -> Any:
    source = presentation.slides[source_index]
    return clone_slide_from(presentation, source)


def clone_slide_from(presentation: Presentation, source: Any) -> Any:
    target = presentation.slides.add_slide(source.slide_layout)
    for shape in source.shapes:
        target.shapes._spTree.insert_element_before(copy.deepcopy(shape.element), "p:extLst")
    return target


def remove_shape(shape: Any) -> None:
    shape.element.getparent().remove(shape.element)


def move_slide_after(presentation: Presentation, slide: Any, after_index: int) -> None:
    slide_ids = presentation.slides._sldIdLst
    slide_id = slide_ids[-1]
    slide_ids.remove(slide_id)
    slide_ids.insert(after_index + 1, slide_id)


def move_slide_before(presentation: Presentation, slide: Any, before_index: int) -> None:
    slide_ids = presentation.slides._sldIdLst
    slide_id = slide_ids[-1]
    slide_ids.remove(slide_id)
    slide_ids.insert(before_index, slide_id)


def expand_team_table(shape: Any, member_count: int, slide_height: int) -> None:
    """Extend the native PowerPoint table instead of moving extra team members outside it."""
    table = shape.table
    wanted_rows = max(2, member_count)
    existing_rows = len(table.rows) - 1
    if wanted_rows <= existing_rows:
        return
    row_height = table.rows[len(table.rows) - 1].height
    for _ in range(wanted_rows - existing_rows):
        table._tbl.append(copy.deepcopy(table.rows[len(table.rows) - 1]._tr))
    maximum_height = slide_height - shape.top - Inches(0.55)
    shape.height = min(shape.height + row_height * (wanted_rows - existing_rows), maximum_height)


def image_slots(count: int) -> list[tuple[float, float, float, float]]:
    if count == 1:
        return [(0.55, 1.15, 8.90, 3.65)]
    if count == 2:
        return [(0.55, 1.15, 4.25, 3.65), (5.15, 1.15, 4.25, 3.65)]
    if count == 3:
        return [(0.55, 1.15, 5.35, 3.65), (6.15, 1.15, 3.25, 1.70), (6.15, 3.10, 3.25, 1.70)]
    return [(0.55, 1.15, 4.25, 1.70), (5.15, 1.15, 4.25, 1.70), (0.55, 3.10, 4.25, 1.70), (5.15, 3.10, 4.25, 1.70)]


PREVENTION_KEYWORDS: dict[str, tuple[str, ...]] = {
    "程序Procedure": ("程序", "流程文件", "process", "procedure", "ecn"),
    "工作指示Work instruction": ("工作指示", "作业指导", "wi"),
    "操作指示SOP": ("sop", "操作指示", "作业规范", "检验"),
    "流程图Flow chart": ("流程图", "flow chart"),
    "失效模式分析D/P-FMEA ": ("fmea", "失效模式"),
    "控制计划Control Plan": ("控制计划", "control plan"),
    "设计规范Design disciplines": ("设计规范", "设计文件", "图纸"),
    "经验教训Lessons Learned": ("经验教训", "教训", "lessons"),
}


def match_prevention_to_d5(countermeasures: list[dict[str, Any]]) -> dict[str, tuple[str, str]]:
    """Match each D7 row's pre-defined item against D5 long-term actions.

    The lookup keys preserve any trailing whitespace that templates may carry,
    so callers should look up with the same raw cell text. A single D5 action
    can satisfy multiple D7 rows; each D7 row only fires once.
    """
    result: dict[str, tuple[str, str]] = {}
    if not countermeasures:
        for item in PREVENTION_KEYWORDS.keys():
            result[item] = ("否", "NA")
        return result
    for d7_item, keywords in PREVENTION_KEYWORDS.items():
        hit_index = -1
        for index, action_item in enumerate(countermeasures):
            action_text = str(action_item.get("action", "")).lower()
            if any(keyword.lower() in action_text for keyword in keywords):
                hit_index = index
                break
        if hit_index >= 0:
            action_text = str(countermeasures[hit_index].get("action", "")).strip()
            result[d7_item] = ("是", f"对策{hit_index + 1}：{action_text}")
        else:
            result[d7_item] = ("否", "NA")
    return result


def _lookup_prevention(inference: dict[str, tuple[str, str]], cell_text: str) -> tuple[str, str]:
    """Match D7 cell text against inference; tolerate trailing whitespace."""
    if cell_text in inference:
        return inference[cell_text]
    stripped = cell_text.strip()
    if stripped in inference:
        return inference[stripped]
    for key, value in inference.items():
        if key.strip() == stripped:
            return value
    return ("否", "NA")


def group_prefix(name: str) -> str | None:
    """Return the leading numeric group prefix for names like '1-1' or '2-3.png'."""
    stem = Path(name).stem
    match = re.match(r"^(\d+)\s*[-_]", stem)
    if not match:
        return None
    return match.group(1)


def plan_image_pages(images: list[Path], plans: list[dict[str, Any]], default_summary: str) -> list[dict[str, Any]]:
    """Force N-M image groups onto one page, then honor the remaining model plan."""
    by_name = {image.name: image for image in images}
    used: set[str] = set()
    result: list[dict[str, Any]] = []
    maximum = {"single": 1, "two": 2, "three": 3, "four": 4}
    fallback_summary = str(default_summary or "").strip() or "相关现场证据见图。"
    summary_by_image: dict[str, str] = {}
    for plan in plans:
        summary = str(plan.get("summary") or "").strip() or fallback_summary
        for name in plan.get("images", []):
            summary_by_image.setdefault(name, summary)

    # A user-provided N-M name is an explicit pagination instruction.  Do this
    # before reading the model plan so a model's "single" layout cannot split
    # 1-1 and 1-2 across slides.
    prefix_buckets: dict[str, list[Path]] = {}
    for image in images:
        prefix = group_prefix(image.name)
        if prefix is not None:
            prefix_buckets.setdefault(prefix, []).append(image)
    for prefix in sorted(prefix_buckets, key=lambda value: int(value)):
        bucket = prefix_buckets[prefix]
        if len(bucket) > 4:
            raise ValueError(f"图片组 {prefix}- 包含 {len(bucket)} 张图片；每页最多支持 4 张。")
        used.update(image.name for image in bucket)
        image_summaries = [summary_by_image.get(image.name, fallback_summary) for image in bucket]
        result.append({"images": bucket, "title": "证据材料", "summary": "；".join(dict.fromkeys(image_summaries)), "image_summaries": image_summaries})

    for plan in plans:
        layout = str(plan.get("layout", "single")).lower()
        selected = [by_name[name] for name in plan.get("images", []) if name in by_name and name not in used]
        if not selected:
            continue
        limit = maximum.get(layout, 1)
        while selected:
            group, selected = selected[:limit], selected[limit:]
            used.update(image.name for image in group)
            summary = text(plan.get("summary") or fallback_summary)
            result.append({"images": group, "title": text(plan.get("title")), "summary": summary, "image_summaries": [summary] * len(group)})

    for image in images:
        if image.name not in used:
            summary = summary_by_image.get(image.name, fallback_summary)
            result.append({"images": [image], "title": "证据材料", "summary": summary, "image_summaries": [summary]})

    return result


def usable_records(records: list[dict[str, Any]], required_field: str) -> list[dict[str, Any]]:
    """Discard blank/model-placeholder records instead of printing fake table data."""
    return [item for item in records if str(item.get(required_field) or "").strip() not in ("", "待确认", "NA")]


def insert_evidence_pages(presentation: Presentation, source_index: int, stage: str, pages: list[dict[str, Any]], before_source: bool = False, images: list[Path] | None = None) -> None:
    insertion_index = source_index
    source = presentation.slides[source_index]
    image_name_set = {image.name for image in (images or [])}
    seen_summaries: list[str] = []
    for page_number, page in enumerate(pages, start=1):
        slide = clone_slide_from(presentation, source)
        for shape in list(slide.shapes):
            if shape.has_table or shape.name == "AI_DYNAMIC" or (shape.has_text_frame and "{{Content}}" in shape.text):
                remove_shape(shape)
        image_summaries = page.get("image_summaries", [])
        if len(page["images"]) > 1 and len(image_summaries) == len(page["images"]):
            body = "；".join(f"图{index}：{summary[:60]}" for index, summary in enumerate(image_summaries, start=1))
        else:
            body = page["summary"].strip()[:200]
        if body.startswith("小结"):
            caption = body
        else:
            caption = f"小结：{body}" if body else "小结："
        if image_name_set:
            for filename in image_name_set:
                caption = caption.replace(filename, "")
        caption = re.sub(r"\s{2,}", " ", caption).strip()
        if not caption.startswith("小结"):
            caption = f"小结：{caption}" if caption else "小结："
        if len(pages) > 1:
            normalized = caption.strip()
            if normalized in seen_summaries:
                caption = f"{caption}（第{page_number}页）"
        seen_summaries.append(caption.strip())
        add_text(slide, caption, 0.55, 4.82 if len(page["images"]) > 1 else 5.03, 8.90, 0.55 if len(page["images"]) > 1 else 0.25, size=8 if len(page["images"]) > 1 else 10)
        for image, slot in zip(page["images"], image_slots(len(page["images"]))):
            crop_picture(slide, image, *slot)
        if before_source:
            move_slide_before(presentation, slide, insertion_index)
        else:
            move_slide_after(presentation, slide, insertion_index)
        insertion_index += 1


def fill_template(template: Path, output: Path, data: dict[str, Any], materials: dict[str, dict[str, Any]]) -> None:
    presentation = Presentation(str(template))
    pages = locate_template_slides(presentation)

    cover, d1 = data["cover"], data["d1"]
    slide = presentation.slides[pages["cover"]]
    replace_content_placeholders(slide, [
        cover.get("prepared_by"), cover.get("reviewed_by"),
        cover.get("approved_by"), cover.get("date"),
    ])
    # The title still supports {{Content}}, while cover approvals can also be ordinary
    # labels such as "编制：" after a customer has customized the template.
    for shape in slide.shapes:
        if shape.has_text_frame and "{{Content}}" in shape.text:
            shape.name = "AI_COVER_TITLE"
        replace_tokens(shape, {"{{Content}}": cover.get("title", "")})
        if not shape.has_text_frame:
            continue
        label_values = (("编制", cover.get("prepared_by")), ("审核", cover.get("reviewed_by")),
                        ("批准", cover.get("approved_by")), ("日期", cover.get("date")))
        for paragraph in shape.text_frame.paragraphs:
            source = paragraph.text.replace(" ", "")
            for label, value in label_values:
                if label in source:
                    display = "" if value in (None, "", "待确认") else str(value).strip()
                    separator = "：" if "：" in paragraph.text or ":" not in paragraph.text else ":"
                    set_paragraph_text(paragraph, f"{label}{separator}{display}")
                    shape.name = "AI_DYNAMIC"
                    break

    # D0 问题了解：按模板中五个 {{content}} 直接写入 AI 返回字段。
    d0 = data.get("d0", {})
    replace_content_placeholders(presentation.slides[pages["d0"]], [
        d0.get("fault_date"), d0.get("model"), d0.get("trace_code"),
        d0.get("station"), d0.get("description"),
    ])

    # D1 团队成立：保留模板的"序号 / 角色"列，其他列由 AI 填入。
    team_shape = first_table(presentation.slides[pages["d1"]])
    expand_team_table(team_shape, len(d1.get("team", [])), presentation.slide_height)
    team_table = team_shape.table
    for row, member in enumerate(d1.get("team", []), start=1):
        put_cell(team_table, row, 1, member.get("name"))
        put_cell(team_table, row, 3, member.get("duty"))
        put_cell(team_table, row, 4, member.get("contact"))
        put_cell(team_table, row, 5, member.get("module"))

    # D2 故障描述：模板保留空白，等待 K3 接入或人工补录。

    # D3 临时措施：模板保留空白，由用户后续手动填写。
    # 故意不写 slide 5 的任何内容。

    # Fill the two D4 5WHY tables from model-generated occurrence and escape reason chains.
    d4 = data.get("d4", {})
    def reason_chain(kind: str) -> list[tuple[str, str]]:
        """Return the problem-and-cause pair for every individual WHY row."""
        direct = d4.get(f"{kind}_why", {})
        problem = str(direct.get("problem", "")).strip()
        values = [str(value).strip() for value in direct.get("whys", []) if str(value or "").strip()]
        rows = d4.get(f"{kind}_why_rows", [])
        if rows:
            return [
                (str(row.get("problem", "")).strip(), str(row.get("cause", "")).strip())
                for row in rows
            ]
        if not values:
            values = [str(item.get("answer", "")).strip() for item in d4.get("five_whys", []) if str(item.get("answer", "")).strip()]
        return [(problem if index == 0 else "", value) for index, value in enumerate(values)]

    occurrence_table = first_table(presentation.slides[pages["d4_occurrence"]]).table
    for row_index in range(1, len(occurrence_table.rows)):
        put_cell(occurrence_table, row_index, 1, "")
        put_cell(occurrence_table, row_index, 2, "")
    occurrence_rows = reason_chain("occurrence")
    for row_index, (problem, cause) in enumerate(occurrence_rows[:len(occurrence_table.rows) - 1], start=1):
        put_cell(occurrence_table, row_index, 1, problem)
        put_cell(occurrence_table, row_index, 2, cause)
    escape_table = first_table(presentation.slides[pages["d4_escape"]]).table
    for row_index in range(1, len(escape_table.rows)):
        put_cell(escape_table, row_index, 1, "")
        put_cell(escape_table, row_index, 2, "")
    escape_rows = reason_chain("escape")
    for row_index, (problem, cause) in enumerate(escape_rows[:len(escape_table.rows) - 1], start=1):
        put_cell(escape_table, row_index, 1, problem)
        put_cell(escape_table, row_index, 2, cause)

    d5_shape = first_table(presentation.slides[pages["d5"]])
    d5_countermeasures = usable_records(data["d5"].get("countermeasures", []), "action")
    expand_team_table(d5_shape, len(d5_countermeasures), presentation.slide_height)
    d5_table = d5_shape.table
    for row_index in range(1, len(d5_table.rows)):
        for col_index in range(len(d5_table.columns)):
            put_cell(d5_table, row_index, col_index, "")
    for row_index, item in enumerate(d5_countermeasures, start=1):
        put_cell(d5_table, row_index, 0, row_index)
        put_cell(d5_table, row_index, 1, item.get("action"))
        put_cell(d5_table, row_index, 2, item.get("owner"))
        put_cell(d5_table, row_index, 3, item.get("date"))
        put_cell(d5_table, row_index, 4, item.get("status") or "进行中")

    d6_shape = first_table(presentation.slides[pages["d6"]])
    validations = usable_records(data["d6"].get("validations", []), "method")
    expand_team_table(d6_shape, len(validations), presentation.slide_height)
    d6_table = d6_shape.table
    for row_index in range(1, len(d6_table.rows)):
        for col_index in range(len(d6_table.columns)):
            put_cell(d6_table, row_index, col_index, "")
    for row_index, item in enumerate(validations, start=1):
        put_cell(d6_table, row_index, 0, row_index)
        put_cell(d6_table, row_index, 1, item.get("method"))
        put_cell(d6_table, row_index, 2, item.get("result"))
        put_cell(d6_table, row_index, 3, item.get("owner"))
        put_cell(d6_table, row_index, 4, item.get("date"))

    # D7 预防措施：优先用模型生成的 D7 数据，无数据时回退关键词匹配。
    d7_data = data.get("d7", {})
    d7_table = first_table(presentation.slides[pages["d7"]]).table
    model_preventions = d7_data.get("preventions", []) if isinstance(d7_data, dict) else []
    if len(model_preventions or []) >= 1 and any(
        isinstance(p, dict) and str(p.get("verdict") or "").strip() == "是"
        for p in (model_preventions or [])
    ):
        prevention_map: dict[str, tuple[str, str]] = {
            str(p.get("item", "")).strip(): (
                str(p.get("verdict", "")).strip() or "否",
                str(p.get("comment", "")).strip() or "NA",
            )
            for p in (model_preventions or [])
            if isinstance(p, dict)
        }
    else:
        prevention_map = match_prevention_to_d5(d5_countermeasures) if d5_countermeasures else {}
    for row in range(1, len(d7_table.rows)):
        verdict, comment = _lookup_prevention(prevention_map, d7_table.cell(row, 0).text)
        put_cell(d7_table, row, 1, verdict)
        put_cell(d7_table, row, 2, comment)
        # owner / date / status 留空待人工填写。

    # D8 结案总结：四个独立 {{content}} 直接对应 AI 的四个结案字段。
    d8 = data.get("d8", {})
    d8_values = [
        d8.get("background"), d8.get("current_status"),
        d8.get("system_strategy"), d8.get("strategy_logic"),
    ]
    # Backward compatibility for saved drafts created before the four-field
    # D8 template: preserve the model's conclusion instead of silently
    # dropping it when the new fields are absent.
    if not any(str(value or "").strip() for value in d8_values):
        d8_values[0] = d8.get("conclusion")
    replace_content_placeholders(presentation.slides[pages["d8"]], [
        *d8_values,
    ])

    # Insert stage evidence immediately after its own slide. Descending order keeps source indexes stable.
    stage_sources = (
        ("D8 结案总结", "d8", pages["d8"]),
        ("D7 预防措施", "d7", pages["d7"]),
        ("D6 效果验证", "d6", pages["d6"]),
        ("D5 长期对策", "d5", pages["d5"]),
        ("D4 原因分析", "d4", pages["d4_evidence"]),
        ("D2 问题描述", "d2", pages["d2"]),
    )
    page_data = data.get("image_pages", {})
    batch_evidence = data.get("_image_evidence", {})
    for stage, key, source_index in stage_sources:
        images = materials[stage]["images"]
        if images:
            stage_pages = list(page_data.get(key, []))
            ev = batch_evidence.get(stage, [])
            if ev and stage_pages:
                evidence_map: dict[str, str] = {s["image"]: s["summary"] for s in ev}
                for p in stage_pages:
                    page_images = list(p.get("images", []))
                    image_summaries = [evidence_map.get(img, "") or p.get("summary", "") for img in page_images]
                    if any(image_summaries) and len(set(image_summaries)) > 1:
                        p["summary"] = "；".join(image_summaries)
                        p["image_summaries"] = image_summaries
            pages = plan_image_pages(images, stage_pages, data.get(key, {}).get("evidence_summary", ""))
            insert_evidence_pages(presentation, source_index, stage, pages, images=images)

    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output))


def main() -> int:
    parser = argparse.ArgumentParser(description="Adayo 8D AI 报告生成器")
    parser.add_argument("--template", type=Path, help="8D 模板 .pptx")
    parser.add_argument("--materials", type=Path, help="含 D0-D8 子文件夹的材料根目录")
    parser.add_argument("--output", type=Path, help="输出 .pptx；省略时写入 EXE 同级目录")
    parser.add_argument("--draft-json", type=Path, help="跳过模型调用，使用已有 JSON 草稿生成 PPT")
    parser.add_argument("--clear", action="store_true", help="直接清空当前目录下 D0-D8 文件夹中的文件")
    parser.add_argument("--yes", action="store_true", help="与 --clear 配合使用，跳过清空确认")
    args = parser.parse_args()
    auto_mode = not any((args.template, args.materials, args.output, args.draft_json, args.clear))
    try:
        if args.clear:
            root = executable_folder()
            if not args.yes and not confirm_clear(root):
                return 0
            n = clear_folders(root)
            print(f"已清空 {n} 个文件。")
            pause()
            return 0

        # Double-click mode: the EXE is placed directly in the material folder.
        if auto_mode:
            while True:
                choice = interactive_menu()
                if choice == 0:
                    return 0
                if choice == 2:
                    if not confirm_clear(executable_folder()):
                        continue
                    n = clear_folders(executable_folder())
                    print(f"已清空 {n} 个文件。")
                    pause()
                    return 0
                # choice == 1: generate PPT
                break
            args.template = bundled_template()
            args.materials = executable_folder()
            args.output = executable_folder() / "8D报告_AI生成.pptx"
        else:
            if not args.template or not args.materials:
                parser.error("必须同时提供 --template 和 --materials。")
            if args.output is None:
                args.output = executable_folder() / "8D报告_AI生成.pptx"

        if not args.template.exists():
            raise FileNotFoundError(f"内置模板不存在：{args.template}")
        print("\n========== 8D 报告生成进度 ==========")
        print("[1/5] 正在读取 00-08（或 D0-D8）材料...")
        materials = collect_materials(args.materials)
        found = sum(len(item["images"]) for item in materials.values())
        if found == 0:
            raise FileNotFoundError("未找到 00-08 或 D0-D8 文件夹中的图片材料。")
        stage_counts = ", ".join(f"{stage[:2]}={len(item['images'])}" for stage, item in materials.items())
        print(f"      已发现 {found} 张图片：{stage_counts}")
        if args.draft_json:
            print("[2/5] 正在读取本地 AI 草稿（跳过模型生成）...")
            data = json.loads(args.draft_json.read_text(encoding="utf-8"))
        else:
            print("[2/5] 正在请求模型生成 8D 内容（模型会自动重试，最长等待 90 秒）...")
            data = generate_content(materials, get_api_key())
            print("[3/5] 模型生成成功，正在保存 AI 草稿...")
            args.output.with_suffix(".json").write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
        warnings = content_warnings(data)
        if warnings:
            print("提示：" + "；".join(warnings) + "。已保留空白字段供人工补充。")
        print("[4/5] 正在写入表格、文字并插入图片，请稍候...")
        fill_template(args.template, args.output, data, materials)
        try:
            print("[5/5] 正在基于完整 PPT 生成报告标题...")
            title = generate_title_from_ppt(args.output, get_api_key())
            if title:
                write_cover_title(args.output, title)
                print(f"      已写入标题：{title}")
        except Exception as exc:
            print(f"提示：PPT 已生成，但标题生成失败：{exc}")
        print(f"完成：已生成可编辑 PPT：{args.output.resolve()}")
        print("====================================")
        return 0
    except Exception as exc:
        print(f"生成失败：{exc}")
        return 1
    finally:
        if auto_mode or getattr(sys, "frozen", False):
            pause()


if __name__ == "__main__":
    sys.exit(main())
